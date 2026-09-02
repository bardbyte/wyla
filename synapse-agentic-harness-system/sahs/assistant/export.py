"""Server-side artifact export (Synapse v2 §5/§13.5): the PPTX deck.

A dashboard exports as a deck — one panel per slide, the meridian
line in the slide NOTES (the disclosure travels with the file, where
a presenter actually looks), the build id in every footer, and an
EXPLORATORY watermark drawn on any slide whose panel still carries
one. A single chart/kpi/table exports as a one-panel deck the same
way. Charts become native PowerPoint charts (editable, not screen-
shots); scatter falls back to a line with markers and says so in the
notes rather than pretending.
"""

from __future__ import annotations

import io
from typing import Any

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)
_CHART_TYPES = {"line": XL_CHART_TYPE.LINE,
                "area": XL_CHART_TYPE.AREA,
                "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
                "scatter": XL_CHART_TYPE.LINE_MARKERS}


def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _text(slide, left, top, width, height, text, *, size=14,
          bold=False, align=PP_ALIGN.LEFT, color=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    para = frame.paragraphs[0]
    para.text = text
    para.alignment = align
    para.font.size = Pt(size)
    para.font.bold = bold
    if color:
        para.font.color.rgb = RGBColor(*color)
    return box


def _footer(slide, build_id: str) -> None:
    _text(slide, Inches(0.4), SLIDE_H - Inches(0.45),
          SLIDE_W - Inches(0.8), Inches(0.35),
          f"Meridian build {build_id or '?'} · every number carries "
          "its definition status", size=9, color=(122, 117, 106))


def _notes(slide, spec: dict[str, Any], extra: str = "") -> None:
    prov = spec.get("provenance") or {}
    lines = []
    if prov:
        lines.append(f"[{prov.get('status', '?')}] "
                     + str(prov.get("meridian_line", "")))
    if spec.get("watermark"):
        lines.append(f"{spec['watermark']}: a check must stand "
                     "behind this number before it loses the "
                     "watermark.")
    if extra:
        lines.append(extra)
    slide.notes_slide.notes_text_frame.text = "\n".join(lines)


def _watermark(slide, text: str) -> None:
    box = _text(slide, Inches(2.5), Inches(3.0), Inches(8.3),
                Inches(1.2), text, size=54, bold=True,
                align=PP_ALIGN.CENTER, color=(200, 170, 110))
    box.rotation = -15


def _chart_slide(slide, spec: dict[str, Any]) -> str:
    kind = spec.get("kind", "line")
    data = CategoryChartData()
    series = spec.get("series") or []
    first = (series[0].get("points") if series else []) or []
    data.categories = [str(p[0]) for p in first]
    for entry in series:
        data.add_series(entry.get("name", "series"),
                        [p[1] for p in entry.get("points", [])])
    slide.shapes.add_chart(
        _CHART_TYPES.get(kind, XL_CHART_TYPE.LINE),
        Inches(0.6), Inches(1.1), SLIDE_W - Inches(1.2),
        Inches(5.4), data)
    return ("scatter rendered as a marked line: PowerPoint charts "
            "here are categorical" if kind == "scatter" else "")


def _table_slide(slide, spec: dict[str, Any]) -> None:
    columns = spec.get("columns") or []
    rows = (spec.get("rows") or [])[:14]
    shape = slide.shapes.add_table(
        len(rows) + 1, max(len(columns), 1), Inches(0.6),
        Inches(1.2), SLIDE_W - Inches(1.2),
        Inches(0.4) * (len(rows) + 1))
    table = shape.table
    for c, col in enumerate(columns):
        table.cell(0, c).text = str(col.get("label", col["key"]))
    for r, row in enumerate(rows, start=1):
        for c, col in enumerate(columns):
            table.cell(r, c).text = str(row.get(col["key"], ""))


def _kpi_slide(slide, spec: dict[str, Any]) -> None:
    if spec.get("label"):
        _text(slide, Inches(1), Inches(2.0), SLIDE_W - Inches(2),
              Inches(0.6), str(spec["label"]).upper(), size=18,
              align=PP_ALIGN.CENTER, color=(122, 117, 106))
    value = f"{spec.get('value', '—')}"
    if spec.get("unit"):
        value += f" {spec['unit']}"
    _text(slide, Inches(1), Inches(2.7), SLIDE_W - Inches(2),
          Inches(1.6), value, size=66, bold=True,
          align=PP_ALIGN.CENTER)
    if isinstance(spec.get("delta"), (int, float)):
        arrow = "▲" if spec["delta"] >= 0 else "▼"
        _text(slide, Inches(1), Inches(4.3), SLIDE_W - Inches(2),
              Inches(0.6), f"{arrow} {abs(spec['delta'])}", size=20,
              align=PP_ALIGN.CENTER,
              color=(46, 125, 50) if spec["delta"] >= 0
              else (179, 57, 47))


def _doc_slide(slide, spec: dict[str, Any]) -> None:
    text = str(spec.get("markdown", ""))
    plain = "\n".join(line.lstrip("#>- ").rstrip()
                      for line in text.splitlines())[:1800]
    _text(slide, Inches(0.6), Inches(1.2), SLIDE_W - Inches(1.2),
          Inches(5.4), plain, size=14)


def _panel_slide(prs: Presentation, ptype: str, title: str,
                 spec: dict[str, Any], build_id: str) -> None:
    slide = _blank(prs)
    _text(slide, Inches(0.4), Inches(0.25), SLIDE_W - Inches(0.8),
          Inches(0.7), title or ptype, size=24, bold=True)
    extra = ""
    if ptype == "chart":
        extra = _chart_slide(slide, spec)
    elif ptype == "table":
        _table_slide(slide, spec)
    elif ptype == "kpi":
        _kpi_slide(slide, spec)
    else:
        _doc_slide(slide, spec)
    if spec.get("watermark"):
        _watermark(slide, str(spec["watermark"]))
    _notes(slide, spec, extra)
    _footer(slide, build_id)


def artifact_pptx(row: dict[str, Any]) -> bytes:
    """A validated artifact row → a .pptx deck, provenance riding in
    the notes. Dashboards get one slide per panel; chart / table /
    kpi / document get a one-panel deck."""
    spec = row.get("spec") or {}
    build_id = str(spec.get("build_id", ""))
    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H

    cover = _blank(prs)
    _text(cover, Inches(0.8), Inches(2.6), SLIDE_W - Inches(1.6),
          Inches(1.2), row.get("title") or row.get("type", "artifact"),
          size=40, bold=True)
    _text(cover, Inches(0.8), Inches(3.9), SLIDE_W - Inches(1.6),
          Inches(0.8),
          f"Synapse · Meridian build {build_id or '?'} · "
          f"v{row.get('version', 1)} · statuses and meridian lines "
          "are in each slide's notes", size=14,
          color=(122, 117, 106))
    _footer(cover, build_id)

    if row.get("type") == "dashboard":
        for panel in spec.get("panels") or []:
            _panel_slide(prs, panel.get("type", "document"),
                         panel.get("title", ""),
                         panel.get("spec") or {}, build_id)
        if spec.get("notes"):
            slide = _blank(prs)
            _text(slide, Inches(0.4), Inches(0.25),
                  SLIDE_W - Inches(0.8), Inches(0.7), "Notes",
                  size=24, bold=True)
            _doc_slide(slide, {"markdown": spec["notes"]})
            _footer(slide, build_id)
    else:
        _panel_slide(prs, row.get("type", "document"),
                     row.get("title", ""), spec, build_id)

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


__all__ = ["artifact_pptx"]
